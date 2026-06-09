#!/usr/bin/env python3
"""Plain PPT: full feature/task list (no status labels), no images/animation."""
from pptx import Presentation
from pptx.util import Pt, Inches

FINAL_DELIVERY = "26 June 2026"

sections = [
    ("Public Website", [
        "Home page", "About Us page", "Contact Us page",
        "Supplier listing page", "Supplier details page",
        "Product listing page (advanced filter + currency)",
        "Product details page (currency)", "FAQs (dynamic)",
        "Content pages (Privacy, Cookie, Terms)",
        "Authentication pages (Login, Registration, Forgot / Reset)",
        "Category system (parent / child / dynamic)",
        "Filter system (country / MOQ / budget)",
        "Unlock supplier (contact details) flow",
        "Signup with Google SSO",
        "Signup with phone verification (OTP)",
        "Country-based pricing (currency by location)",
    ]),
    ("Buyer Features", [
        "Dashboard", "Post RFQs", "Manage RFQs", "Manage Quotations",
        "Profile settings (update profile, change password)",
        "Notification settings (Email / SMS / WhatsApp / Telegram)",
        "Token / Credit purchase + payment + usage list",
        "Save data (saved suppliers / RFQs / favorite products)",
        "Switch between Buyer / Supplier",
        "Update phone with verification",
    ]),
    ("Supplier Features", [
        "Dashboard", "Listing & browsing RFQs + RFQ details",
        "RFQ responses / Manage Quotations",
        "Product management (list / view / edit / delete)",
        "Company settings (profile, documents, payment & shipping, strength)",
        "Verification (document submission)",
        "Profile settings (update profile, change password)",
        "Notification settings (Email / SMS / WhatsApp / Telegram)",
        "Switch between Buyer / Supplier",
        "Update phone with verification",
    ]),
    ("Admin Features", [
        "Login / Forgot / Reset password", "Dashboard",
        "User management (view / approve / reject / block)",
        "Impersonate login (login as buyer / supplier)",
        "RFQ management (view all / delete spam)",
        "Verification (view documents / approve / reject)",
        "Verified badge (manual)",
        "Product management",
        "Category management (add / edit / delete)",
        "Content management (Privacy / Terms / Cookie / pages)",
        "FAQ management",
        "Credit / token management (add / remove / transactions)",
        "Trust & Safety: report / flag users",
        "Trust & Safety: spam protection (limit RFQs)",
        "Internal roles & permissions management",
        "Country management", "Country port management",
        "Notification settings (enable / disable Telegram / WhatsApp)",
        "Settings (update profile, change password)",
        "Audit logging",
    ]),
]

technology = [
    ("Framework", "Next.js / React", "SEO-friendly server rendering; full-stack app"),
    ("Language", "TypeScript", "Type safety across the application"),
    ("Database", "PostgreSQL", "Relational data (users, RFQs, quotes, credits)"),
    ("ORM", "Prisma", "Type-safe database access & migrations"),
    ("Authentication", "Auth.js (NextAuth)", "Email, Google SSO, roles, impersonation"),
    ("Validation", "Zod", "Client & server input validation"),
    ("UI / Styling", "Tailwind CSS + shadcn/ui", "Modern, accessible interface"),
    ("Payments", "Stripe", "Secure credit-pack checkout"),
    ("Email", "Resend", "Transactional email"),
    ("SMS / WhatsApp", "Twilio", "Phone verification & alerts"),
    ("Messaging", "Telegram Bot API", "Opt-in notifications"),
    ("File Storage", "Cloudflare R2 / S3", "Documents & media"),
    ("Hosting", "Any Node.js host", "Railway / Render / DigitalOcean / AWS / VPS"),
    ("Database Hosting", "Managed PostgreSQL", "Scalable managed database"),
]

prs = Presentation()
TITLE_SLIDE = prs.slide_layouts[0]
TITLE_CONTENT = prs.slide_layouts[1]
TITLE_ONLY = prs.slide_layouts[5]


def bullet_slide(title, bullets, size=16):
    slide = prs.slides.add_slide(TITLE_CONTENT)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(size)


def feature_slides(title, items, limit=13):
    if len(items) <= limit:
        bullet_slide(title, items)
    else:
        half = (len(items) + 1) // 2
        bullet_slide(title + " (1/2)", items[:half])
        bullet_slide(title + " (2/2)", items[half:])


# 1. Title
s = prs.slides.add_slide(TITLE_SLIDE)
s.shapes.title.text = "B2B Marketplace Platform"
s.placeholders[1].text = "Proposed Development Plan"

# 2. Overview
bullet_slide("Overview", [
    "This document presents the features proposed for development.",
    "All items listed below are planned for development.",
    "Users: Buyer, Supplier, Admin",
    "Modules: Public Website, Buyer, Supplier, Admin",
    "Proposed delivery date: " + FINAL_DELIVERY,
], size=18)

# 3+. Proposed feature lists
for title, items in sections:
    feature_slides("Planned Features: " + title, items)

# Technology stack (table)
slide = prs.slides.add_slide(TITLE_ONLY)
slide.shapes.title.text = "Proposed Technology Stack"
rows = len(technology) + 1
table = slide.shapes.add_table(rows, 3, Inches(0.5), Inches(1.4), Inches(9.0), Inches(0.3 * rows)).table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(2.8)
table.columns[2].width = Inches(4.0)
for c, h in enumerate(["Area", "Technology", "Purpose"]):
    cell = table.cell(0, c); cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
for r, (a, t, p) in enumerate(technology, 1):
    for c, val in enumerate((a, t, p)):
        cell = table.cell(r, c); cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(11)

# Proposed delivery
s = prs.slides.add_slide(TITLE_SLIDE)
s.shapes.title.text = "Proposed Delivery Date"
s.placeholders[1].text = FINAL_DELIVERY

out = "B2B_Documentation.pptx"
prs.save(out)
print("Saved:", out, "| Slides:", len(prs.slides._sldIdLst))
